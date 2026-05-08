"""
Generate a 12-slide project presentation with real outputs and code snippets.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


PROJECT_ROOT = Path(__file__).resolve().parents[1]
OUTPUTS_DIR = PROJECT_ROOT / "outputs"
DELIVERABLES_DIR = PROJECT_ROOT / "deliverables"


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, text):
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(9.2), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = RGBColor(10, 40, 70)

    line = slide.shapes.add_connector(1, Inches(0.4), Inches(0.95), Inches(9.6), Inches(0.95))
    line.line.color.rgb = RGBColor(230, 90, 70)
    line.line.width = Pt(2.5)


def add_bullets(slide, items, left=0.5, top=1.15, width=9.0, height=5.9):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i > 0:
            tf.add_paragraph()
        p = tf.paragraphs[i]

        text, level, bold = item
        p.text = text
        p.level = level
        p.font.bold = bold
        p.font.size = Pt(20 - (level * 2))
        p.font.color.rgb = RGBColor(45, 45, 45)
        p.space_after = Pt(5)


def add_image(slide, image_path, left, top, width):
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width))
    else:
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.9))
        p = box.text_frame.paragraphs[0]
        p.text = f"Missing image: {image_path.name}"
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(180, 20, 20)


def add_code_box(slide, code_text, left, top, width, height):
    rect = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(28, 34, 40)
    rect.line.color.rgb = RGBColor(70, 80, 90)

    tf = rect.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code_text
    p.font.name = "Consolas"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(230, 240, 250)


def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, RGBColor(243, 248, 252))

    box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9.0), Inches(1.7))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "Quantum Graph Convolutional Network"
    p.font.size = Pt(46)
    p.font.bold = True
    p.font.color.rgb = RGBColor(8, 42, 75)
    p.alignment = PP_ALIGN.CENTER

    tf.add_paragraph()
    p2 = tf.paragraphs[1]
    p2.text = "Power Grid Reliability Prediction"
    p2.font.size = Pt(30)
    p2.font.color.rgb = RGBColor(230, 90, 70)
    p2.alignment = PP_ALIGN.CENTER

    tf.add_paragraph()
    p3 = tf.paragraphs[2]
    p3.text = "Final Project Presentation (12 Slides)"
    p3.font.size = Pt(18)
    p3.font.color.rgb = RGBColor(70, 70, 70)
    p3.alignment = PP_ALIGN.CENTER


def slide_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, RGBColor(248, 251, 253))
    add_title(slide, "1. Problem Statement and Objective")
    add_bullets(slide, [
        ("Goal: predict unstable grid states before cascading failures.", 0, True),
        ("Grid modeled as a graph: substations = nodes, lines = edges.", 0, False),
        ("Need to combine node signals + network connectivity patterns.", 0, False),
        ("Deliverable target: reproducible CPU run with zero training epochs.", 0, False),
        ("Project scope: phase-wise implementation, outputs, baselines, documents.", 0, False),
    ])


def slide_pipeline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, RGBColor(248, 251, 253))
    add_title(slide, "2. End-to-End Pipeline")
    add_bullets(slide, [
        ("Phase 1: build IEEE 14-bus graph + synthetic operational data.", 0, True),
        ("Phase 2: encode node features into quantum states.", 0, True),
        ("Phase 3: variational quantum circuit + graph message passing.", 0, True),
        ("Phase 4: hybrid readout, metrics, visual outputs.", 0, True),
        ("Execution mode used: CPU-only, num_epochs = 0 (no training loop).", 0, False),
    ])


def slide_phase1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, RGBColor(248, 251, 253))
    add_title(slide, "3. Phase 1 - Graph Modeling (Code + Output)")

    code = (
        "# phase1_data_modeling.py\n"
        "grid = PowerGridSimulator(config=GRID_CONFIG)\n"
        "X, y = grid.generate_node_features(num_samples=250)\n"
        "grid.visualize_grid(save_path='outputs/01_grid_topology.png')"
    )
    add_code_box(slide, code, left=0.5, top=1.25, width=4.4, height=2.0)
    add_image(slide, OUTPUTS_DIR / "01_grid_topology.png", left=5.1, top=1.25, width=4.3)

    add_bullets(slide, [
        ("Output proof: IEEE 14-bus topology image generated.", 0, False),
        ("Dataset shape: (samples, 14 nodes, 5 features).", 0, False),
    ], left=0.5, top=3.45, width=9.0, height=2.8)


def slide_phase2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, RGBColor(248, 251, 253))
    add_title(slide, "4. Phase 2 - Quantum Feature Embedding")

    code = (
        "# phase2_quantum_embedding.py\n"
        "for each qubit i:\n"
        "    qml.RX(theta_i, wires=i)\n"
        "    qml.RY(phi_i, wires=i)\n"
        "entangle with CNOT ladder\n"
        "measure [expval(PauliZ(i))]"
    )
    add_code_box(slide, code, left=0.5, top=1.25, width=4.4, height=2.35)
    add_image(slide, OUTPUTS_DIR / "05_quantum_states.png", left=5.1, top=1.25, width=4.3)

    add_bullets(slide, [
        ("Fix applied: removed constant embedding collapse from old encoding.", 0, False),
        ("Now each node gets informative 4-value quantum representation.", 0, False),
    ], left=0.5, top=3.85, width=9.0, height=2.5)


def slide_phase3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, RGBColor(248, 251, 253))
    add_title(slide, "5. Phase 3 - QGCN Layer and Message Passing")

    code = (
        "# phase3_quantum_layer.py\n"
        "node_output = VQC(node_embedding)\n"
        "neighbor_msg = mean(neighbor_embeddings)\n"
        "updated = 0.7 * node_output + 0.3 * neighbor_msg\n"
        "classical = dot(updated, output_weights) + bias"
    )
    add_code_box(slide, code, left=0.5, top=1.25, width=4.6, height=2.2)

    add_bullets(slide, [
        ("Quantum part extracts non-linear node features.", 0, False),
        ("Graph part propagates neighborhood context.", 0, False),
        ("Hybrid readout returns per-node risk scores.", 0, False),
    ], left=0.5, top=3.65, width=4.6, height=2.6)

    add_image(slide, OUTPUTS_DIR / "04_grid_risk_heatmap.png", left=5.3, top=1.25, width=4.1)


def slide_phase4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, RGBColor(248, 251, 253))
    add_title(slide, "6. Phase 4 - CPU-Friendly Inference Mode")

    code = (
        "# config.py\n"
        "TRAINING_CONFIG = {\n"
        "    'num_epochs': 0,\n"
        "    ...\n"
        "}\n"
        "# phase4_training.py\n"
        "if epochs <= 0: return"
    )
    add_code_box(slide, code, left=0.5, top=1.25, width=4.2, height=2.2)
    add_image(slide, OUTPUTS_DIR / "02_training_history.png", left=4.9, top=1.25, width=4.5)

    add_bullets(slide, [
        ("No GPU needed: full pipeline runs with zero epochs.", 0, False),
        ("Still executes real quantum forward path for demonstration.", 0, False),
    ], left=0.5, top=3.65, width=9.0, height=2.2)


def slide_outputs_a(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, RGBColor(248, 251, 253))
    add_title(slide, "7. Output Demonstration A")
    add_image(slide, OUTPUTS_DIR / "03_prediction_distribution.png", left=0.5, top=1.2, width=4.5)
    add_image(slide, OUTPUTS_DIR / "06_roc_curve.png", left=5.0, top=1.2, width=4.5)

    add_bullets(slide, [
        ("Left: probability distribution for normal vs failed states.", 0, False),
        ("Right: ROC plot generated from inference outputs.", 0, False),
    ], left=0.5, top=5.35, width=9.0, height=1.4)


def slide_outputs_b(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, RGBColor(248, 251, 253))
    add_title(slide, "8. Output Demonstration B")
    add_image(slide, OUTPUTS_DIR / "04_grid_risk_heatmap.png", left=0.5, top=1.2, width=4.5)
    add_image(slide, OUTPUTS_DIR / "07_baseline_comparison.png", left=5.0, top=1.2, width=4.5)

    add_bullets(slide, [
        ("Heatmap highlights relative node-level risk in the grid.", 0, False),
        ("Baseline chart compares voltage-rule and logistic models.", 0, False),
    ], left=0.5, top=5.35, width=9.0, height=1.4)


def slide_metrics(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, RGBColor(248, 251, 253))
    add_title(slide, "9. Run Results (Current CPU Execution)")

    add_bullets(slide, [
        ("QGCN (zero-epoch inference):", 0, True),
        ("Accuracy 0.1400 | Precision 0.1400 | Recall 1.0000", 1, False),
        ("F1 0.2456 | AUC 0.5615", 1, False),
        ("", 0, False),
        ("Classical baselines:", 0, True),
        ("Voltage Rule -> Accuracy 0.880, Recall 0.400, F1 0.571", 1, False),
        ("Logistic Regression -> Accuracy 0.680, Recall 0.100, F1 0.111", 1, False),
        ("", 0, False),
        ("Reference file: docs/EXECUTION_SUMMARY.md", 0, False),
    ], left=0.6, top=1.2, width=8.8, height=5.9)


def slide_workdone(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, RGBColor(248, 251, 253))
    add_title(slide, "10. Work Completed and Fixes")

    add_bullets(slide, [
        ("Implemented full 4-phase project structure.", 0, True),
        ("Fixed quantum embedding collapse in angle encoding.", 0, True),
        ("Made pipeline stable with no-epoch CPU mode.", 0, True),
        ("Replaced synthetic demo with real QGCN inference outputs.", 0, True),
        ("Updated reports/presentation to remove fabricated static metrics.", 0, True),
        ("Verified one-click deliverables script end-to-end.", 0, True),
    ])


def slide_files(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, RGBColor(248, 251, 253))
    add_title(slide, "11. Deliverables and Demo Commands")

    add_bullets(slide, [
        ("Main deliverables:", 0, True),
        ("deliverables/QGCN_PowerGrid_Presentation_12_Slides.pptx", 1, False),
        ("deliverables/QGCN_Project_Writeup_Udai_Ratinam_G.docx", 1, False),
        ("docs/EXECUTION_SUMMARY.md", 1, False),
        ("outputs/01..07 PNG figures", 1, False),
        ("", 0, False),
        ("Demo commands:", 0, True),
        ("python src/demo.py", 1, False),
        ("python src/baseline_comparison.py", 1, False),
        ("powershell -ExecutionPolicy Bypass -File ./run_deliverables.ps1", 1, False),
    ])


def slide_conclusion(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, RGBColor(243, 248, 252))
    add_title(slide, "12. Conclusion and Next Steps")

    add_bullets(slide, [
        ("Project is complete, reproducible, and CPU-friendly.", 0, True),
        ("Quantum implementation is now demonstrable from real code paths.", 0, True),
        ("Outputs and reports are regenerated and submission-ready.", 0, True),
        ("Next: optional non-zero epoch experiments on stronger hardware.", 0, False),
        ("Next: larger IEEE grids and stronger ablation studies.", 0, False),
    ], left=0.6, top=1.4, width=8.8, height=4.5)

    thank_you = slide.shapes.add_textbox(Inches(0.5), Inches(6.1), Inches(9.0), Inches(0.7))
    p = thank_you.text_frame.paragraphs[0]
    p.text = "Thank You"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(230, 90, 70)


def build_presentation() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_title(prs)
    slide_problem(prs)
    slide_pipeline(prs)
    slide_phase1(prs)
    slide_phase2(prs)
    slide_phase3(prs)
    slide_phase4(prs)
    slide_outputs_a(prs)
    slide_outputs_b(prs)
    slide_metrics(prs)
    slide_workdone(prs)
    slide_files(prs)
    slide_conclusion(prs)

    DELIVERABLES_DIR.mkdir(parents=True, exist_ok=True)
    output_path = DELIVERABLES_DIR / "QGCN_PowerGrid_Presentation_12_Slides.pptx"
    try:
        prs.save(str(output_path))
        return output_path
    except PermissionError:
        # File may be open in PowerPoint; save a timestamped copy instead
        from datetime import datetime
        alt_path = DELIVERABLES_DIR / (
            f"QGCN_PowerGrid_Presentation_12_Slides_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        )
        prs.save(str(alt_path))
        print(f"Note: could not overwrite existing file. Saved as: {alt_path}")
        return alt_path
    except Exception as exc:
        print(f"Failed to save presentation: {exc}")
        raise


if __name__ == "__main__":
    ppt = build_presentation()
    print(f"Created: {ppt}")
