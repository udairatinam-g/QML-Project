"""
Generate PowerPoint presentation for QGCN Project
Theoretical Understanding for Academic Audiences
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 32, 72)  # Dark blue
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(78, 205, 196)  # Teal
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_items, bg_color=RGBColor(240, 245, 250)):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(15, 32, 72)  # Dark blue
    
    # Add underline
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(1.15), Inches(9.5), Inches(1.15))
    line.line.color.rgb = RGBColor(255, 107, 107)  # Red
    line.line.width = Pt(3)
    
    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        
        if isinstance(item, tuple):
            # (text, level, is_bold)
            text, level, is_bold = item
            p.text = text
            p.level = level
            p.font.size = Pt(18 - level * 2)
            p.font.bold = is_bold
        else:
            p.text = item
            p.level = 0
            p.font.size = Pt(18)
        
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_before = Pt(6)
        p.space_after = Pt(6)
    
    return slide

def add_image_slide(prs, title, image_path):
    """Add a slide with title and image"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 245, 250)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(15, 32, 72)
    
    # Add underline
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(1.15), Inches(9.5), Inches(1.15))
    line.line.color.rgb = RGBColor(255, 107, 107)
    line.line.width = Pt(3)
    
    # Add image
    try:
        slide.shapes.add_picture(image_path, Inches(1), Inches(1.5), width=Inches(8))
    except:
        # If image doesn't exist, add placeholder
        text_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(3))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"[Image: {image_path}]"
        p.font.size = Pt(18)
        p.alignment = PP_ALIGN.CENTER
    
    return slide

# ============================================================================
# SLIDE 1: Title Slide
# ============================================================================
add_title_slide(prs, 
    "Quantum Graph Convolutional Networks",
    "A Theory-to-Implementation Guide for Power Grid Reliability")

# ============================================================================
# SLIDE 2: Abstract
# ============================================================================
add_content_slide(prs, "Abstract", [
    ("This project develops a Quantum Graph Convolutional Network (QGCN) to predict power grid risk.", 0, False),
    ("The power grid is modeled as a network of connected substations.", 1, False),
    ("The model learns from voltage, power, and load patterns across the full network.", 1, False),
    ("", 0, False),
    ("Main Contribution:", 0, True),
    ("A complete phase-wise implementation from data modeling to risk prediction.", 1, False),
    ("A hybrid approach that combines graph learning ideas with quantum feature processing.", 1, False),
    ("", 0, False),
    ("Outcome:", 0, True),
    ("The system provides an early warning score for possible grid instability.", 1, False),
    ("This helps operators take preventive action and improve reliability.", 1, False),
])

# ============================================================================
# SLIDE 3: Problem Statement & Motivation
# ============================================================================
add_content_slide(prs, "Problem Statement & Motivation", [
    ("Goal: Predict grid instability early so operators can act before failure", 0, True),
    ("Power grid = graph (substations are nodes, transmission lines are edges)", 1, False),
    ("Failures spread through connections, not isolated points", 1, False),
    ("", 0, False),
    ("Why Classical Methods Struggle:", 0, True),
    ("Need to model both node features and network topology", 1, False),
    ("Long-range dependencies are hard to capture efficiently", 1, False),
    ("High-dimensional state estimation is computationally expensive", 1, False),
    ("", 0, False),
    ("QGCN Idea (Easy View):", 0, True),
    ("Encode each node into qubits", 1, False),
    ("Entangle neighboring nodes to share information", 1, False),
    ("Read out one risk score per node and one global risk score", 1, False),
])

# ============================================================================
# SLIDE 4: Simple Theory Behind the Model
# ============================================================================
add_content_slide(prs, "Simple Theory Behind the Model", [
    ("Each substation has operational features such as voltage, power, and load.", 0, False),
    ("These features are converted into quantum-friendly values for processing.", 1, False),
    ("This conversion helps the model capture complex relationships efficiently.", 1, False),
    ("", 0, False),
    ("Why Quantum Processing Helps:", 0, True),
    ("It can represent many interaction patterns in a compact way.", 1, False),
    ("It captures non-linear behavior that is common in real power systems.", 1, False),
    ("It supports network-level reasoning instead of isolated node decisions.", 1, False),
    ("", 0, False),
    ("Practical Meaning:", 0, True),
    ("The model remains data-driven and practical for engineering use.", 1, False),
    ("Better feature representation leads to earlier and more reliable warnings.", 1, False),
])

# ============================================================================
# SLIDE 5: Model Architecture in Simple Steps
# ============================================================================
add_content_slide(prs, "Model Architecture in Simple Steps", [
    ("Step 1: Feature Encoding", 0, True),
    ("Convert each substation's features into a quantum-ready format.", 1, False),
    ("", 0, False),
    ("Step 2: Quantum Processing", 0, True),
    ("Use trainable quantum layers to learn important signal patterns.", 1, False),
    ("", 0, False),
    ("Step 3: Message Passing", 0, True),
    ("Share information between connected substations.", 1, False),
    ("This mimics how disturbances spread through the grid.", 1, False),
    ("", 0, False),
    ("Step 4: Risk Prediction", 0, True),
    ("Convert learned node patterns into a clear probability of instability.", 1, False),
    ("Use this probability as an early-warning indicator for operators.", 1, False),
])

# ============================================================================
# SLIDE 6: Implementation Blueprint (Phase-by-Phase)
# ============================================================================
add_content_slide(prs, "Implementation Blueprint (Phase-by-Phase)", [
    ("Phase 1 - Data Modeling (phase1_data_modeling.py)", 0, True),
    ("Create IEEE 14-bus graph; generate X shape (N, 14, 5) and labels y", 1, False),
    ("", 0, False),
    ("Phase 2 - Quantum Embedding (phase2_quantum_embedding.py)", 0, True),
    ("Convert each node's 5 features into 4-qubit encoded representation", 1, False),
    ("", 0, False),
    ("Phase 3 - Quantum Message Passing (phase3_quantum_layer.py)", 0, True),
    ("Apply VQC + neighbor interaction to produce robust node embeddings", 1, False),
    ("", 0, False),
    ("Phase 4 - Training (phase4_training.py)", 0, True),
    ("Train the model to separate stable and unstable operating conditions.", 1, False),
    ("Output: probability of instability for each scenario", 1, False),
])

# ============================================================================
# SLIDE 7: Results, Impact, and Easy Interpretation
# ============================================================================
add_content_slide(prs, "How It Helps: Results in Easy Terms", [
    ("Model Performance (IEEE 14-bus, 250 samples):", 0, True),
    ("Accuracy 82.9%, Precision 79.4%, Recall 81.6%, F1 80.5%, AUC 0.847", 1, False),
    ("Interpretation: catches most risky states while keeping false alarms manageable", 1, False),
    ("", 0, False),
    ("Operational Value:", 0, True),
    ("Early-warning signal for dispatch/control-room teams", 1, False),
    ("Supports preventive actions: load balancing, rerouting, reserve activation", 1, False),
    ("Can be integrated as a risk-scoring module in smart-grid software", 1, False),
    ("", 0, False),
    ("Simple Mental Model:", 0, True),
    ("Input: grid snapshot -> Process: quantum graph reasoning -> Output: risk probability", 1, False),
    ("If probability is high, operator gets an actionable alert", 1, False),
    ("Result: better reliability, lower outage risk, safer smart-city infrastructure", 1, False),
])

# ============================================================================
# SLIDE 8: Baseline Comparison
# ============================================================================
add_content_slide(prs, "Baseline Comparison (Why QGCN Is Useful)", [
    ("We compared QGCN against simple classical baselines.", 0, False),
    ("", 0, False),
    ("Voltage Rule Baseline:", 0, True),
    ("High accuracy on easy cases, but low recall for failure detection.", 1, False),
    ("Accuracy 0.88, Recall 0.40, F1 0.57", 1, False),
    ("", 0, False),
    ("Logistic Regression Baseline:", 0, True),
    ("Struggles on complex network interactions.", 1, False),
    ("Accuracy 0.68, Recall 0.10, F1 0.11", 1, False),
    ("", 0, False),
    ("QGCN Project Result:", 0, True),
    ("Balanced detection for real operations.", 1, False),
    ("Accuracy 0.829, Recall 0.816, F1 0.805", 1, False),
    ("Conclusion: QGCN is better when the priority is catching risky states early.", 1, False),
])

# Save presentation
deliverables_dir = Path("deliverables")
deliverables_dir.mkdir(parents=True, exist_ok=True)
output_path = str(deliverables_dir / "QGCN_PowerGrid_Presentation.pptx")
try:
    prs.save(output_path)
except PermissionError:
    output_path = str(deliverables_dir / "QGCN_PowerGrid_Presentation_Simplified.pptx")
    prs.save(output_path)

print(f"✓ PowerPoint presentation created: {output_path}")
print("\nSlide Summary:")
print("  Slide 1: Title Slide")
print("  Slide 2: Abstract")
print("  Slide 3: Problem Statement & Motivation")
print("  Slide 4: Simple Theory Behind the Model")
print("  Slide 5: Model Architecture in Simple Steps")
print("  Slide 6: Implementation Blueprint")
print("  Slide 7: Results, Impact, and Easy Interpretation")
print("  Slide 8: Baseline Comparison")
